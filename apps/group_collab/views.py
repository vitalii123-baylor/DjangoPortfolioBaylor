import base64
import hashlib
import json
import math
from django.shortcuts import render, get_object_or_404, redirect
from django.http import JsonResponse
from django.views.decorators.http import require_POST, require_GET
from django.conf import settings
from google import genai

from .models import StudyGroup, Member, Vote, Presentation, PresentationPart

_client = None

def _get_client():
    global _client
    if _client is None:
        _client = genai.Client(api_key=settings.GEMINI_API_KEY)
    return _client

def _ensure_session(request):
    if not request.session.session_key:
        request.session.create()
    return request.session.session_key

def _get_member(request, group):
    sid = request.session.session_key
    if not sid:
        return None
    return group.members.filter(session_key=sid).first()

def _voter_hash(session_key, group_code):
    return hashlib.sha256(f"{session_key}::{group_code}".encode()).hexdigest()

def _vote_counts(group):
    counts = {}
    for v in group.votes.all():
        counts[v.candidate_id] = counts.get(v.candidate_id, 0) + 1
    return counts

def _slides_per_part(duration_min, n):
    if not duration_min:
        return 3
    per = max(1, duration_min / max(1, n))
    return max(2, min(6, round(per / 2)))


def _build_parts_template(n, spp):
    """Global slide numbering: Part 1 → slides 1..spp, Part 2 → spp+1..2*spp, etc."""
    parts = []
    for i in range(1, n + 1):
        start = (i - 1) * spp + 1
        slides = '\n\n'.join(
            f'[SLIDE {start + j}] <slide {start + j} title>\n<speaker notes for slide {start + j}>'
            for j in range(spp)
        )
        parts.append(f'[PART {i}]\n{slides}')
    return '\n\n'.join(parts)


def _parse_slides_in_part(part_text):
    import re
    slides = []
    matches = list(re.finditer(r'\[SLIDE (\d+)\]\s*(.+?)(?:\n|$)', part_text))
    for j, m in enumerate(matches):
        num   = int(m.group(1))          # keep the real global number
        title = m.group(2).strip()
        start = m.end()
        end   = matches[j + 1].start() if j + 1 < len(matches) else len(part_text)
        slides.append({'number': num, 'title': title, 'text': part_text[start:end].strip()})
    return slides


def _parse_ai_response(text, n):
    import re
    full_text = ''
    parts = []

    if 'FULL_TEXT:' in text:
        fs = text.index('FULL_TEXT:') + len('FULL_TEXT:')
        pm = text.find('PARTS:')
        full_text = text[fs: pm if pm != -1 else len(text)].strip()

    for i in range(1, n + 1):
        marker = f'[PART {i}]'
        nxt = f'[PART {i + 1}]' if i < n else None
        if marker in text:
            s = text.index(marker) + len(marker)
            e = text.index(nxt) if nxt and nxt in text else len(text)
            raw = text[s:e].strip()
            slides = _parse_slides_in_part(raw)
            parts.append({'text': raw, 'slides': slides})
        else:
            parts.append({'text': '', 'slides': []})

    if not full_text and parts:
        full_text = '\n\n'.join(p['text'] for p in parts)
    return full_text, parts


# ── Landing ──────────────────────────────────────────────────────────────────

def group_home(request):
    return render(request, 'group_collab/home.html')


def _duration_hint(duration_min, n):
    if not duration_min:
        return ''
    per = duration_min / n
    words = int(per * 130)
    return (
        f"\nTotal duration: {duration_min} minutes ({per:.1f} min per presenter)."
        f" Each part should be ~{words} words (speaking pace ~130 words/minute)."
    )


@require_POST
def demo_generate(request):
    data = json.loads(request.body)
    topic = data.get('topic', '').strip()
    try:
        n = max(2, min(int(data.get('n', 3)), 6))
    except (ValueError, TypeError):
        n = 3
    try:
        duration = int(data.get('duration', 0))
    except (ValueError, TypeError):
        duration = 0

    if not topic:
        return JsonResponse({'ok': False, 'error': 'Topic is required'})

    spp = _slides_per_part(duration, n)
    parts_tmpl = _build_parts_template(n, spp)
    prompt = f"""You are helping a student group prepare a presentation.

Topic: {topic}
Number of presenters: {n}{_duration_hint(duration, n)}
Slides per presenter: {spp}

Write a clear, engaging, well-structured presentation script on this topic.
Divide it into exactly {n} parts. Each part has exactly {spp} slides with a title and speaker notes.

Reply in this EXACT format (no text before FULL_TEXT:):

FULL_TEXT:
<complete script here>

PARTS:
{parts_tmpl}"""

    try:
        resp = _get_client().models.generate_content(model='gemini-2.5-flash-lite', contents=prompt)
        full_text, parts = _parse_ai_response(resp.text, n)
        return JsonResponse({
            'ok': True,
            'full_text': full_text,
            'parts': [{'part_number': i + 1, 'text': p['text'], 'slides': p['slides']}
                      for i, p in enumerate(parts)],
        })
    except Exception as e:
        err = str(e)
        if '429' in err or 'quota' in err.lower():
            return JsonResponse({'ok': False, 'error': 'AI quota exceeded — try again later'})
        return JsonResponse({'ok': False, 'error': err})


@require_POST
def demo_pdf(request):
    upload = request.FILES.get('pdf')
    if not upload:
        return JsonResponse({'ok': False, 'error': 'Please upload a file'})
    fname = upload.name.lower()
    if not fname.endswith(('.pdf', '.pptx')):
        return JsonResponse({'ok': False, 'error': 'Please upload a .pdf or .pptx file'})
    try:
        n = max(2, min(int(request.POST.get('n', 3)), 6))
    except (ValueError, TypeError):
        n = 3
    try:
        duration = int(request.POST.get('duration', 0))
    except (ValueError, TypeError):
        duration = 0

    page_texts = []
    page_images = []   # one base64 data-URL (or None) per page/slide

    if fname.endswith('.pdf'):
        try:
            import fitz  # pymupdf
            file_bytes = upload.read()
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            mat = fitz.Matrix(1.67, 1.67)  # ~43 DPI — compact thumbnails
            for page in doc:
                page_texts.append(page.get_text() or '')
                pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
                img_bytes = pix.tobytes(output="jpeg", jpg_quality=88)
                page_images.append(f"data:image/jpeg;base64,{base64.b64encode(img_bytes).decode()}")
            doc.close()
        except ImportError:
            return JsonResponse({'ok': False, 'error': 'pymupdf not installed — run: pip install pymupdf'})
        except Exception as e:
            return JsonResponse({'ok': False, 'error': f'PDF error: {e}'})
    else:
        try:
            from pptx import Presentation as PptxPres
            prs = PptxPres(upload)
            for slide in prs.slides:
                texts = []
                img_data = None
                for shape in slide.shapes:
                    if img_data is None and getattr(shape, 'shape_type', None) == 13:
                        try:
                            blob = shape.image.blob
                            if len(blob) <= 800_000:
                                ct = shape.image.content_type or 'image/png'
                                img_data = f"data:{ct};base64,{base64.b64encode(blob).decode()}"
                        except Exception:
                            pass
                    if hasattr(shape, 'text') and shape.text.strip():
                        texts.append(shape.text.strip())
                page_texts.append('\n'.join(texts))
                page_images.append(img_data)
        except ImportError:
            return JsonResponse({'ok': False, 'error': 'python-pptx not installed'})
        except Exception as e:
            return JsonResponse({'ok': False, 'error': f'PPTX read error: {e}'})

    num_pages = len(page_texts)
    raw_text = '\n'.join(page_texts).strip()
    if not raw_text:
        return JsonResponse({'ok': False, 'error': 'Could not extract text from this file'})
    raw_text = raw_text[:14000]

    spp = max(1, math.ceil(num_pages / n))
    if duration:
        spp = max(spp, _slides_per_part(duration, n))
    total_slides = n * spp
    parts_tmpl = _build_parts_template(n, spp)
    file_label = 'PPTX slides' if fname.endswith('.pptx') else 'PDF pages'
    prompt = f"""You are helping a student group create a presentation from their uploaded file.
Number of presenters: {n}{_duration_hint(duration, n)}
{file_label}: {num_pages} — distribute content evenly, {spp} slides per presenter ({total_slides} slides total, numbered {1}–{total_slides} continuously across ALL parts).

Content:
\"\"\"
{raw_text}
\"\"\"

1. Extract key information and restructure it into a clear, engaging presentation script.
2. Remove page numbers, headers, footnotes, and formatting artifacts.
3. Divide into exactly {n} parts, each with exactly {spp} slides (title + speaker notes).
4. Slide numbers must be CONTINUOUS across all parts — do NOT restart from 1 in each part.

Reply in this EXACT format:

FULL_TEXT:
<complete restructured script>

PARTS:
{parts_tmpl}"""

    try:
        resp = _get_client().models.generate_content(model='gemini-2.5-flash-lite', contents=prompt)
        full_text, parts = _parse_ai_response(resp.text, n)
        return JsonResponse({
            'ok': True,
            'full_text': full_text,
            'parts': [
                {
                    'part_number': i + 1,
                    'text': p['text'],
                    'slides': p['slides'],
                    'page_start': i * spp + 1,
                    'page_end': min((i + 1) * spp, num_pages),
                    'page_images': page_images[i * spp: (i + 1) * spp],
                }
                for i, p in enumerate(parts)
            ],
        })
    except Exception as e:
        err = str(e)
        if '429' in err or 'quota' in err.lower():
            return JsonResponse({'ok': False, 'error': 'AI quota exceeded — try again later'})
        return JsonResponse({'ok': False, 'error': err})


@require_POST
def demo_extract(request):
    """Extract PDF or PPTX content per page/slide, split evenly across N speakers."""
    upload = request.FILES.get('file')
    if not upload:
        return JsonResponse({'ok': False, 'error': 'No file uploaded'})
    try:
        n = max(2, min(int(request.POST.get('n', 3)), 6))
    except (ValueError, TypeError):
        n = 3

    name = upload.name.lower()
    raw_slides = []

    if name.endswith('.pdf'):
        try:
            import fitz  # pymupdf
            file_bytes = upload.read()
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            mat = fitz.Matrix(1.67, 1.67)
            for i, page in enumerate(doc):
                text = page.get_text().strip()
                pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
                img_data = f"data:image/jpeg;base64,{base64.b64encode(pix.tobytes(output='jpeg', jpg_quality=88)).decode()}"
                raw_slides.append({
                    'number': i + 1,
                    'title': f'Page {i + 1}',
                    'text': text or '(no text extracted)',
                    'image': img_data,
                })
            doc.close()
        except ImportError:
            return JsonResponse({'ok': False, 'error': 'pymupdf not installed — run: pip install pymupdf'})
        except Exception as e:
            return JsonResponse({'ok': False, 'error': f'PDF error: {e}'})

    elif name.endswith('.pptx'):
        try:
            from pptx import Presentation
            prs = Presentation(upload)
            for i, slide in enumerate(prs.slides):
                title = ''
                bodies = []
                img_data = None
                for shape in slide.shapes:
                    if img_data is None and getattr(shape, 'shape_type', None) == 13:
                        try:
                            blob = shape.image.blob
                            if len(blob) <= 800_000:
                                ct = shape.image.content_type or 'image/png'
                                img_data = f"data:{ct};base64,{base64.b64encode(blob).decode()}"
                        except Exception:
                            pass
                    if not hasattr(shape, 'text'):
                        continue
                    t = shape.text.strip()
                    if not t:
                        continue
                    ph = getattr(shape, 'placeholder_format', None)
                    if ph and ph.idx == 0:
                        title = t
                    else:
                        bodies.append(t)
                if not title and bodies:
                    title = bodies.pop(0)
                raw_slides.append({
                    'number': i + 1,
                    'title': title or f'Slide {i + 1}',
                    'text': '\n'.join(bodies),
                    'image': img_data,
                })
        except ImportError:
            return JsonResponse({'ok': False, 'error': 'python-pptx not installed — run: pip install python-pptx'})
        except Exception as e:
            return JsonResponse({'ok': False, 'error': f'PPTX error: {e}'})
    else:
        return JsonResponse({'ok': False, 'error': 'Please upload a .pdf or .pptx file'})

    if not raw_slides:
        return JsonResponse({'ok': False, 'error': 'No content found in file'})

    total = len(raw_slides)
    per = math.ceil(total / n)
    parts = []
    for i in range(n):
        chunk = raw_slides[i * per:(i + 1) * per]
        if chunk:
            parts.append({
                'part_number': i + 1,
                'slides': chunk,
                'page_start': chunk[0]['number'],
                'page_end': chunk[-1]['number'],
            })

    return JsonResponse({'ok': True, 'parts': parts, 'total': total, 'filename': upload.name})


@require_POST
def get_tips(request):
    data = json.loads(request.body)
    script = data.get('script', '').strip()[:12000]
    if not script:
        return JsonResponse({'ok': False, 'error': 'No script provided'})

    prompt = f"""You are an expert presentation coach. Analyze this student presentation script and give specific, actionable improvement tips.

SCRIPT:
\"\"\"
{script}
\"\"\"

Provide exactly 6 tips covering: opening hook, structure & flow, language clarity, speaker transitions, closing impact, and delivery/body language.

Format EACH tip EXACTLY like this (no extra text):

[TIP 1] <short title>
<2-3 sentence explanation with a concrete suggestion>

[TIP 2] <short title>
<2-3 sentence explanation with a concrete suggestion>

[TIP 3] <short title>
...

[TIP 4] <short title>
...

[TIP 5] <short title>
...

[TIP 6] <short title>
..."""

    try:
        resp = _get_client().models.generate_content(model='gemini-2.5-flash-lite', contents=prompt)
        tips = _parse_tips(resp.text)
        return JsonResponse({'ok': True, 'tips': tips})
    except Exception as e:
        err = str(e)
        if '429' in err or 'quota' in err.lower():
            return JsonResponse({'ok': False, 'error': 'AI quota exceeded — try again later'})
        return JsonResponse({'ok': False, 'error': err})


@require_POST
def rephrase_part(request):
    data = json.loads(request.body)
    part_text = data.get('text', '').strip()[:8000]
    feedback = data.get('feedback', '').strip()[:500]

    if not part_text:
        return JsonResponse({'ok': False, 'error': 'No text provided'})

    fb_line = f'\nUser feedback: "{feedback}"' if feedback else '\nMake it clearer, more engaging, and better structured.'

    prompt = f"""You are a presentation coach. Rewrite this presentation section based on the feedback.

CURRENT TEXT:
\"\"\"
{part_text}
\"\"\"
{fb_line}

Rewrite it keeping the slide structure. Use 2–4 slides.

Format EXACTLY like this (no extra text):

[SLIDE 1] <slide title>
<speaker notes — what to say>

[SLIDE 2] <slide title>
<speaker notes>

[SLIDE 3] <slide title>
<speaker notes>"""

    try:
        resp = _get_client().models.generate_content(model='gemini-2.5-flash-lite', contents=prompt)
        new_text = resp.text.strip()
        slides = _parse_slides_in_part(new_text)
        return JsonResponse({'ok': True, 'text': new_text, 'slides': slides})
    except Exception as e:
        err = str(e)
        if '429' in err or 'quota' in err.lower():
            return JsonResponse({'ok': False, 'error': 'AI quota exceeded'})
        return JsonResponse({'ok': False, 'error': err})


def _parse_tips(text):
    import re
    tips = []
    matches = re.findall(r'\[TIP \d+\]\s*(.+?)\n([\s\S]+?)(?=\[TIP \d+\]|$)', text)
    for title, body in matches:
        tips.append({'title': title.strip(), 'body': body.strip()})
    return tips


# ── Create / Join ─────────────────────────────────────────────────────────────

@require_POST
def create_group(request):
    _ensure_session(request)
    name = request.POST.get('name', '').strip()[:100]
    nickname = request.POST.get('nickname', '').strip()[:50]
    try:
        max_members = max(2, min(int(request.POST.get('max_members', 4)), 20))
    except ValueError:
        max_members = 4

    if not name or not nickname:
        return redirect('group_collab:home')

    group = StudyGroup.objects.create(name=name, max_members=max_members)
    Member.objects.create(
        group=group,
        nickname=nickname,
        session_key=request.session.session_key,
        is_leader=True,
        role='leader',
    )
    return redirect('group_collab:dashboard', code=group.code)


@require_POST
def join_group(request):
    _ensure_session(request)
    code = request.POST.get('code', '').strip().upper()
    nickname = request.POST.get('nickname', '').strip()[:50]

    try:
        group = StudyGroup.objects.get(code=code)
    except StudyGroup.DoesNotExist:
        return redirect('group_collab:home')

    existing = group.members.filter(session_key=request.session.session_key).first()
    if existing:
        return redirect('group_collab:dashboard', code=code)

    if group.members.count() >= group.max_members or not nickname:
        return redirect('group_collab:home')

    Member.objects.create(
        group=group,
        nickname=nickname,
        session_key=request.session.session_key,
    )
    return redirect('group_collab:dashboard', code=code)


# ── Dashboard ─────────────────────────────────────────────────────────────────

def group_dashboard(request, code):
    _ensure_session(request)
    group = get_object_or_404(StudyGroup, code=code)
    member = _get_member(request, group)

    if not member:
        return redirect('group_collab:home')

    latest_pres = group.presentations.prefetch_related('parts__member').order_by('-created_at').first()
    voter_hash = _voter_hash(request.session.session_key, code)
    my_vote = None
    if group.voting_open:
        v = group.votes.filter(voter_hash=voter_hash).first()
        if v:
            my_vote = v.candidate_id

    demo_topics = [
        'AI in Modern Healthcare',
        'Climate Change Solutions',
        'Cybersecurity Best Practices',
        'The Future of Remote Work',
        'Blockchain & Decentralized Finance',
    ]
    context = {
        'group': group,
        'me': member,
        'members': group.members.all().order_by('joined_at'),
        'vote_counts': _vote_counts(group),
        'my_vote': my_vote,
        'latest_pres': latest_pres,
        'roles': Member.ROLES,
        'demo_topics': demo_topics,
    }
    return render(request, 'group_collab/dashboard.html', context)


# ── Members management ────────────────────────────────────────────────────────

@require_POST
def update_group(request, code):
    _ensure_session(request)
    group = get_object_or_404(StudyGroup, code=code)
    me = _get_member(request, group)
    if not me or not me.is_leader:
        return JsonResponse({'ok': False})

    data = json.loads(request.body)
    if 'max_members' in data:
        group.max_members = max(group.members.count(), min(int(data['max_members']), 20))
    group.save()
    return JsonResponse({'ok': True, 'max_members': group.max_members})


@require_POST
def kick_member(request, code, member_id):
    _ensure_session(request)
    group = get_object_or_404(StudyGroup, code=code)
    me = _get_member(request, group)
    if not me or not me.is_leader:
        return JsonResponse({'ok': False})

    target = get_object_or_404(Member, id=member_id, group=group)
    if target.id == me.id:
        return JsonResponse({'ok': False, 'error': "Can't remove yourself"})
    target.delete()
    return JsonResponse({'ok': True})


@require_POST
def update_role(request, code, member_id):
    _ensure_session(request)
    group = get_object_or_404(StudyGroup, code=code)
    me = _get_member(request, group)
    if not me or not me.is_leader:
        return JsonResponse({'ok': False})

    data = json.loads(request.body)
    role = data.get('role', 'member')
    valid_roles = [r[0] for r in Member.ROLES]
    if role not in valid_roles:
        return JsonResponse({'ok': False})

    Member.objects.filter(id=member_id, group=group).update(role=role)
    return JsonResponse({'ok': True})


def leave_group(request, code):
    _ensure_session(request)
    group = get_object_or_404(StudyGroup, code=code)
    member = _get_member(request, group)
    if member:
        if member.is_leader:
            others = group.members.exclude(id=member.id).order_by('joined_at')
            if others.exists():
                nxt = others.first()
                nxt.is_leader = True
                nxt.role = 'leader'
                nxt.save()
        member.delete()
    return redirect('group_collab:home')


# ── Voting ────────────────────────────────────────────────────────────────────

@require_POST
def open_voting(request, code):
    _ensure_session(request)
    group = get_object_or_404(StudyGroup, code=code)
    me = _get_member(request, group)
    if not me or not me.is_leader:
        return JsonResponse({'ok': False})
    group.votes.all().delete()
    group.voting_open = True
    group.save()
    return JsonResponse({'ok': True})


@require_POST
def cast_vote(request, code):
    _ensure_session(request)
    group = get_object_or_404(StudyGroup, code=code)
    me = _get_member(request, group)

    if not me or not group.voting_open:
        return JsonResponse({'ok': False, 'error': 'Voting is closed'})

    data = json.loads(request.body)
    try:
        candidate = group.members.get(id=int(data.get('candidate_id')))
    except (Member.DoesNotExist, TypeError, ValueError):
        return JsonResponse({'ok': False, 'error': 'Invalid candidate'})

    voter_hash = _voter_hash(request.session.session_key, code)
    Vote.objects.update_or_create(
        group=group, voter_hash=voter_hash,
        defaults={'candidate': candidate},
    )
    return JsonResponse({'ok': True, 'counts': _vote_counts(group)})


@require_POST
def close_voting(request, code):
    _ensure_session(request)
    group = get_object_or_404(StudyGroup, code=code)
    me = _get_member(request, group)
    if not me or not me.is_leader:
        return JsonResponse({'ok': False})

    counts = _vote_counts(group)
    if counts:
        winner_id = max(counts, key=lambda x: counts[x])
        group.members.update(is_leader=False, role='member')
        winner = group.members.get(id=winner_id)
        winner.is_leader = True
        winner.role = 'leader'
        winner.save()

    group.voting_open = False
    group.votes.all().delete()
    group.save()
    return JsonResponse({'ok': True})


# ── Polling ───────────────────────────────────────────────────────────────────

@require_GET
def poll_status(request, code):
    _ensure_session(request)
    group = get_object_or_404(StudyGroup, code=code)
    me = _get_member(request, group)
    if not me:
        return JsonResponse({'ok': False})

    voter_hash = _voter_hash(request.session.session_key, code)
    my_vote_id = None
    if group.voting_open:
        v = group.votes.filter(voter_hash=voter_hash).first()
        if v:
            my_vote_id = v.candidate_id

    return JsonResponse({
        'ok': True,
        'voting_open': group.voting_open,
        'max_members': group.max_members,
        'member_count': group.members.count(),
        'my_vote': my_vote_id,
        'vote_counts': _vote_counts(group),
        'members': [
            {'id': m.id, 'nickname': m.nickname, 'is_leader': m.is_leader, 'role': m.role}
            for m in group.members.all().order_by('joined_at')
        ],
    })


# ── AI Presentation ───────────────────────────────────────────────────────────

@require_POST
def generate_presentation(request, code):
    _ensure_session(request)
    group = get_object_or_404(StudyGroup, code=code)
    me = _get_member(request, group)
    if not me:
        return JsonResponse({'ok': False, 'error': 'Not a member'})

    data = json.loads(request.body)
    topic = data.get('topic', '').strip()
    if not topic:
        return JsonResponse({'ok': False, 'error': 'Topic is required'})
    try:
        duration = int(data.get('duration', 0))
    except (ValueError, TypeError):
        duration = 0

    n = group.members.count()
    spp = _slides_per_part(duration, n)
    parts_tmpl = _build_parts_template(n, spp)

    prompt = f"""You are helping a student group prepare a presentation.

Topic: {topic}
Number of presenters: {n}{_duration_hint(duration, n)}
Total slides: {n * spp} (slides {1}–{n * spp}, numbered continuously across ALL parts — do NOT restart from 1).
Slides per presenter: {spp}

Write a clear, engaging, well-structured presentation script on this topic.
Divide it into exactly {n} parts, each with exactly {spp} slides (title + speaker notes).

Reply in this EXACT format (no text before FULL_TEXT:):

FULL_TEXT:
<complete script here>

PARTS:
{parts_tmpl}"""

    try:
        resp = _get_client().models.generate_content(model='gemini-2.5-flash-lite', contents=prompt)
        full_text, parts = _parse_ai_response(resp.text, n)

        pres = Presentation.objects.create(group=group, topic=topic, full_text=full_text)
        members_list = list(group.members.all().order_by('joined_at'))
        for i, p in enumerate(parts):
            PresentationPart.objects.create(
                presentation=pres,
                member=members_list[i] if i < len(members_list) else None,
                part_number=i + 1,
                text=p['text'],
            )

        return JsonResponse({
            'ok': True,
            'pres_id': pres.id,
            'full_text': full_text,
            'parts': [
                {'id': pt.id, 'part_number': pt.part_number, 'text': pt.text,
                 'slides': _parse_slides_in_part(pt.text),
                 'member': pt.member.nickname if pt.member else f'Part {pt.part_number}'}
                for pt in pres.parts.all()
            ],
        })
    except Exception as e:
        err = str(e)
        if '429' in err or 'quota' in err.lower():
            return JsonResponse({'ok': False, 'error': 'AI quota exceeded — try again later'})
        return JsonResponse({'ok': False, 'error': err})


@require_POST
def upload_pdf(request, code):
    _ensure_session(request)
    group = get_object_or_404(StudyGroup, code=code)
    me = _get_member(request, group)
    if not me:
        return JsonResponse({'ok': False, 'error': 'Not a member'})

    pdf_file = request.FILES.get('pdf')
    if not pdf_file or not pdf_file.name.lower().endswith('.pdf'):
        return JsonResponse({'ok': False, 'error': 'Please upload a .pdf file'})

    try:
        import pypdf
        reader = pypdf.PdfReader(pdf_file)
        num_pages = len(reader.pages)
        raw_text = '\n'.join(page.extract_text() or '' for page in reader.pages).strip()
        if not raw_text:
            return JsonResponse({'ok': False, 'error': 'Could not extract text from PDF'})
        raw_text = raw_text[:14000]
    except ImportError:
        return JsonResponse({'ok': False, 'error': 'pypdf not installed. Run: pip install pypdf'})
    except Exception as e:
        return JsonResponse({'ok': False, 'error': f'PDF read error: {e}'})

    try:
        duration = int(request.POST.get('duration', 0))
    except (ValueError, TypeError):
        duration = 0

    n = group.members.count()
    # Distribute PDF pages fairly across presenters
    spp = max(1, math.ceil(num_pages / n))
    if duration:
        spp = max(spp, _slides_per_part(duration, n))
    total_slides = n * spp
    parts_tmpl = _build_parts_template(n, spp)

    prompt = f"""You are helping a student group create a presentation from their uploaded PDF material.
Number of presenters: {n}{_duration_hint(duration, n)}
PDF pages: {num_pages} — distribute content evenly, {spp} slides per presenter ({total_slides} slides total, numbered {1}–{total_slides} continuously across ALL parts).

PDF content:
\"\"\"
{raw_text}
\"\"\"

1. Extract key information and restructure into a clear, engaging script.
2. Remove page numbers, headers, footnotes, and artifacts.
3. Divide into exactly {n} parts, each with exactly {spp} slides (title + speaker notes).
4. Slide numbers must be CONTINUOUS across all parts — do NOT restart from 1 in each part.

Reply in this EXACT format:

FULL_TEXT:
<complete restructured script>

PARTS:
{parts_tmpl}"""

    try:
        resp = _get_client().models.generate_content(model='gemini-2.5-flash-lite', contents=prompt)
        full_text, parts = _parse_ai_response(resp.text, n)

        pres = Presentation.objects.create(
            group=group, topic=f'PDF: {pdf_file.name}', full_text=full_text
        )
        members_list = list(group.members.all().order_by('joined_at'))
        for i, p in enumerate(parts):
            PresentationPart.objects.create(
                presentation=pres,
                member=members_list[i] if i < len(members_list) else None,
                part_number=i + 1,
                text=p['text'],
            )

        return JsonResponse({
            'ok': True,
            'pres_id': pres.id,
            'full_text': full_text,
            'parts': [
                {'id': pt.id, 'part_number': pt.part_number, 'text': pt.text,
                 'slides': _parse_slides_in_part(pt.text),
                 'member': pt.member.nickname if pt.member else f'Part {pt.part_number}'}
                for pt in pres.parts.all()
            ],
        })
    except Exception as e:
        err = str(e)
        if '429' in err or 'quota' in err.lower():
            return JsonResponse({'ok': False, 'error': 'AI quota exceeded — try again later'})
        return JsonResponse({'ok': False, 'error': err})


@require_POST
def update_notes(request, code, part_id):
    _ensure_session(request)
    group = get_object_or_404(StudyGroup, code=code)
    me = _get_member(request, group)
    if not me:
        return JsonResponse({'ok': False})

    part = get_object_or_404(PresentationPart, id=part_id, presentation__group=group)
    data = json.loads(request.body)
    part.notes = data.get('notes', '')[:1000]
    part.save()
    return JsonResponse({'ok': True})
